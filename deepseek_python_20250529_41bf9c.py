import tkinter as tk
from tkinter import filedialog, messagebox
import pandas as pd
import numpy as np
import matplotlib.pyplot as plt
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
import re
import json
import os
from pathlib import Path

# Configuration file path
CONFIG_FILE = "user_settings.json"

def load_user_settings():
    """Load user settings from JSON file"""
    if os.path.exists(CONFIG_FILE):
        with open(CONFIG_FILE, 'r') as f:
            return json.load(f)
    return {
        "last_pch_path": "",
        "last_pptx_path": "",
        "remember_paths": False
    }

def save_user_settings(settings):
    """Save user settings to JSON file"""
    with open(CONFIG_FILE, 'w') as f:
        json.dump(settings, f)

def create_landscape_presentation():
    """Create a new landscape PowerPoint presentation"""
    prs = Presentation()
    
    # Set slide width to 13.33 inches (landscape) and height to 7.5 inches
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)
    
    # Add a blank slide layout (layout 5 is typically blank)
    blank_slide_layout = prs.slide_layouts[5]
    
    # Add a title slide
    title_slide = prs.slides.add_slide(prs.slide_layouts[0])
    title = title_slide.shapes.title
    subtitle = title_slide.placeholders[1]
    
    title.text = "Vibration Analysis Report"
    subtitle.text = "Generated on " + pd.Timestamp.now().strftime('%Y-%m-%d')
    
    return prs

def get_or_create_presentation(pptx_path=None):
    """Get existing presentation or create new landscape one"""
    if pptx_path and os.path.exists(pptx_path):
        try:
            return Presentation(pptx_path)
        except Exception as e:
            messagebox.showwarning("Template Error", 
                                 f"Could not load template: {e}\nCreating new presentation.")
    
    return create_landscape_presentation()

def extract_data(text):
    """Extract data from PCH file text"""
    point_id_pattern = re.compile(r'\$POINT ID =\s+(\d+)')
    data_pattern = re.compile(r'(\d\.\d+E[+-]\d+)\s+(\d\.\d+E[+-]\d+)\s+(\d\.\d+E[+-]\d+)\s+(\d\.\d+E[+-]\d+)')

    point_ids = point_id_pattern.findall(text)
    sections = point_id_pattern.split(text)[1:]

    data_dict = {point_id: [] for point_id in point_ids}

    for i in range(0, len(sections), 2):
        point_id = sections[i]
        section_data = sections[i + 1]
        matches = data_pattern.findall(section_data)

        for match in matches:
            frequency, x_dir, y_dir, z_dir = match
            data_dict[point_id].append([frequency, x_dir, y_dir, z_dir])

    df_list = []
    for point_id, data in data_dict.items():
        df = pd.DataFrame(data, columns=["Frequency", 
                                        f"Node_{point_id}_tn_x_file_1", 
                                        f"Node_{point_id}_tn_y_file_1", 
                                        f"Node_{point_id}_tn_z_file_1"])
        df_list.append(df)

    return pd.concat(df_list, axis=1)

def update_presentation(start_node, end_node, pptx_path, extracted_df2):
    """Update PowerPoint with vibration analysis plots"""
    # Node title mapping (same as before)
    node_title_map = {
        8000001: "Engine Mount Front Top LH", 
        # ... rest of your node mappings ...
    }

    nodes = set(col.split('_')[1] for col in extracted_df2.columns if col != 'Frequency')
    df_RSS = pd.DataFrame()
    df_RSS['Frequency'] = extracted_df2['Frequency']

    for node in nodes:
        x_col = f'Node_{node}_tn_x_file_1'
        y_col = f'Node_{node}_tn_y_file_1'
        z_col = f'Node_{node}_tn_z_file_1'
        res_col = f'RSS_{node}'
        df_RSS[res_col] = np.sqrt(extracted_df2[x_col]**2 + extracted_df2[y_col]**2 + extracted_df2[z_col]**2)

    clr = ["#1f77b4", "#ff7f0e", "#2ca02c", "#d62728", "#9467bd", 
           "#8c564b", "#e377c2", "#7f7f7f", "#bcbd22", "#17becf"]

    filtered_columns = [col for col in df_RSS.columns if re.search(r'RSS_(\d+)', col) 
                       and start_node <= int(re.search(r'RSS_(\d+)', col).group(1)) <= end_node]

    rms_df = pd.DataFrame({
        "RMS_1-100": np.sqrt((df_RSS.query("0 < Frequency < 100").iloc[:, 1:] ** 2).mean()),
        "RMS_100-150": np.sqrt((df_RSS.query("100 < Frequency < 150").iloc[:, 1:] ** 2).mean()),
        "RMS_150-300": np.sqrt((df_RSS.query("150 < Frequency < 300").iloc[:, 1:] ** 2).mean())
    }).T / 1000

    prs = get_or_create_presentation(pptx_path)
    slide_layout = prs.slide_layouts[5]

    num_plots = len(filtered_columns)
    font_size = max(10, 28 - min(num_plots, 4) * 2)

    plt.rcParams.update({
        'axes.titlesize': font_size,
        'axes.labelsize': font_size,
        'xtick.labelsize': font_size,
        'ytick.labelsize': font_size,
        'figure.autolayout': True
    })

    for i, node in enumerate(filtered_columns):
        if i % 4 == 0:
            slide = prs.slides.add_slide(slide_layout)
            left_margin = Inches(0.5)
            top_margin = Inches(0.5)
            plot_width = Inches(6)
            plot_height = Inches(4.5)

        row, col = divmod(i % 4, 2)
        left = left_margin + col * (plot_width + Inches(0.5))
        top = top_margin + row * (plot_height + Inches(0.5))

        fig, ax = plt.subplots(figsize=(plot_width.inches, plot_height.inches))
        ax = rms_df[node].plot(kind='bar', color=clr[:3], ax=ax, width=0.8)
        
        # Plot formatting (same as before)
        # ...

        plot_path = f'temp_plot_{i}.png'
        fig.savefig(plot_path, dpi=300, bbox_inches='tight')
        plt.close(fig)

        slide.shapes.add_picture(plot_path, left, top, width=plot_width, height=plot_height)
        os.remove(plot_path)

    output_path = "Vibration_Analysis_Report.pptx"
    prs.save(output_path)
    return output_path

def main():
    """Main application flow"""
    settings = load_user_settings()
    
    root = tk.Tk()
    root.withdraw()

    # Ask for PCH file
    pch_path = filedialog.askopenfilename(
        title="Select PCH File",
        filetypes=[("Punch files", "*.pch")],
        initialdir=os.path.dirname(settings.get("last_pch_path", ""))
    )
    
    if not pch_path:
        return

    # Ask for PPTX template (optional)
    pptx_path = filedialog.askopenfilename(
        title="Select PowerPoint Template (Optional)",
        filetypes=[("PowerPoint files", "*.pptx")],
        initialdir=os.path.dirname(settings.get("last_pptx_path", ""))
    )

    # Update settings if user wants to remember paths
    if messagebox.askyesno("Remember Paths", "Remember these paths for next time?"):
        settings.update({
            "last_pch_path": pch_path,
            "last_pptx_path": pptx_path if pptx_path else "",
            "remember_paths": True
        })
        save_user_settings(settings)

    try:
        with open(pch_path, 'r') as f:
            pch_content = f.read()
        
        extracted_df = extract_data(pch_content)
        extracted_df2 = extracted_df.loc[:, ~(extracted_df.columns == 'Frequency') & ~extracted_df.columns.duplicated()].apply(pd.to_numeric, errors='coerce')

        output_pptx = update_presentation(
            start_node=8000001,
            end_node=8000045,
            pptx_path=pptx_path if pptx_path else None,
            extracted_df2=extracted_df2
        )

        messagebox.showinfo("Success", f"Report generated successfully:\n{output_pptx}")
        os.startfile(output_pptx)

    except Exception as e:
        messagebox.showerror("Error", f"An error occurred:\n{str(e)}")

if __name__ == "__main__":
    main()